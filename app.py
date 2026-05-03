import os
import sys
import subprocess
import threading
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import zipfile
from io import BytesIO
import sqlite3
import gc
import time

import pystray
from PIL import Image, ImageDraw
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# ==========================================
# WATCHDOG EVENT HANDLER
# ==========================================
class DocumentEventHandler(FileSystemEventHandler):
    def __init__(self, app_instance):
        self.app = app_instance
        self.allowed_extensions = {".pdf", ".docx", ".pptx", ".xlsx"}

    def on_created(self, event):
        self.process_event(event)

    def on_modified(self, event):
        self.process_event(event)
        
    def process_event(self, event):
        if event.is_directory:
            return
        filepath = event.src_path
        ext = os.path.splitext(filepath)[1].lower()
        if ext in self.allowed_extensions:
            try:
                os_mtime = os.path.getmtime(filepath)
                insert_or_update_document(filepath, "", "Đang chờ", os_mtime)
                
                # Tự động gắp file "Đang chờ"
                if not getattr(self.app, 'is_indexing', False):
                    self.app.is_indexing = True
                    threading.Thread(target=self.app.background_processing_worker, daemon=True).start()
            except Exception as e:
                print(f"Watchdog lỗi xử lý file: {e}")

# ==========================================
# CẤU HÌNH TỐI ƯU CUDNN VÀ AI (LAZY LOADING)
# ==========================================
use_gpu = False

classifier = None
ocr_reader = None
ai_initialized = False

CANDIDATE_LABELS = ["Toán học", "Lập trình", "Kiến thức Đại cương", "Khác"]

def init_ai_models(status_callback=None):
    global classifier, ocr_reader, ai_initialized, use_gpu
    if ai_initialized:
        return
        
    import torch
    from transformers import pipeline
    import easyocr
    
    torch.backends.cudnn.benchmark = True
    use_gpu = torch.cuda.is_available()

    if status_callback:
        status_callback("Trạng thái: Đang nạp mô hình NLP Zero-Shot...")
    print("Đang khởi tạo mô hình NLP Zero-Shot (FP16)... Vui lòng đợi.")
    try:
        classifier = pipeline(
            "zero-shot-classification", 
            model="MoritzLaurer/mDeBERTa-v3-base-mnli-xnli", 
            device=0 if use_gpu else -1, 
            torch_dtype=torch.float16 if use_gpu else torch.float32
        )
    except Exception as e:
        print(f"Lỗi khởi tạo mô hình AI: {e}")
        classifier = None

    if status_callback:
        status_callback("Trạng thái: Đang nạp EasyOCR...")
    try:
        ocr_reader = easyocr.Reader(['vi', 'en'], gpu=use_gpu)
    except Exception as e:
        print(f"Lỗi khởi tạo EasyOCR: {e}")
        ocr_reader = None
        
    ai_initialized = True

# ==========================================
# CƠ SỞ DỮ LIỆU SQLITE
# ==========================================
DB_PATH = 'file_index.db'

def init_db():
    try:
        conn = sqlite3.connect(DB_PATH, timeout=10)
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                filepath TEXT UNIQUE,
                content TEXT,
                topic TEXT,
                last_modified REAL
            )
        ''')
        try:
            cursor.execute("ALTER TABLE documents ADD COLUMN last_modified REAL")
        except sqlite3.OperationalError:
            pass # Cột đã tồn tại
        conn.commit()
        conn.close()
    except Exception as e:
        print(f"Lỗi tạo DB: {e}")

def get_db_mtime(filepath):
    try:
        conn = sqlite3.connect(DB_PATH, timeout=10)
        cursor = conn.cursor()
        cursor.execute("SELECT last_modified FROM documents WHERE filepath = ?", (filepath,))
        row = cursor.fetchone()
        conn.close()
        if row and row[0]:
            return row[0]
        return 0.0
    except:
        return 0.0

def insert_or_update_document(filepath, content, topic, mtime):
    try:
        conn = sqlite3.connect(DB_PATH, timeout=10)
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO documents (filepath, content, topic, last_modified)
            VALUES (?, ?, ?, ?)
            ON CONFLICT(filepath) DO UPDATE SET
            content=excluded.content,
            topic=excluded.topic,
            last_modified=excluded.last_modified
        ''', (filepath, content, topic, mtime))
        conn.commit()
        conn.close()
    except Exception as e:
        print(f"Lỗi lưu DB cho file {filepath}: {e}")

def search_db(keyword, topic):
    try:
        conn = sqlite3.connect(DB_PATH, timeout=10)
        cursor = conn.cursor()
        if topic == "Tất cả":
            cursor.execute('''
                SELECT filepath, content, topic FROM documents
                WHERE content LIKE ? OR filepath LIKE ?
            ''', (f'%{keyword}%', f'%{keyword}%'))
        else:
            cursor.execute('''
                SELECT filepath, content, topic FROM documents
                WHERE (content LIKE ? OR filepath LIKE ?) AND topic = ?
            ''', (f'%{keyword}%', f'%{keyword}%', topic))
        results = cursor.fetchall()
        conn.close()
        return results
    except Exception as e:
        print(f"Lỗi truy vấn DB: {e}")
        return []

# ==========================================
# TRÍCH XUẤT VĂN BẢN (I/O CPU)
# ==========================================
def extract_text_pdf(filepath):
    import PyPDF2
    from pdf2image import convert_from_path, pdfinfo_from_path
    import numpy as np
    import torch
    text = []
    try:
        reader = PyPDF2.PdfReader(filepath)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text.append(extracted)
    except Exception as e:
        pass
        
    combined_text = "\n".join(text)
    if len(combined_text) > 50:
        return combined_text
        
    try:
        current_dir = os.path.dirname(os.path.abspath(__file__))
        poppler_relative_path = os.path.join(current_dir, 'poppler', 'Library', 'bin')
        
        info = pdfinfo_from_path(filepath, poppler_path=poppler_relative_path)
        total_pages = info.get("Pages", 0)
        
        batch_size = 3
        for start_page in range(1, total_pages + 1, batch_size):
            end_page = min(start_page + batch_size - 1, total_pages)
            images = convert_from_path(
                filepath, 
                first_page=start_page, 
                last_page=end_page, 
                poppler_path=poppler_relative_path
            )
            
            for img in images:
                if ocr_reader:
                    try:
                        img_array = np.array(img)
                        ocr_result = ocr_reader.readtext(img_array, detail=0)
                        if ocr_result:
                            text.append("\n".join(ocr_result))
                    except Exception as e:
                        print("Bỏ qua OCR do thiếu thư viện")
            
            del images
            if use_gpu:
                torch.cuda.empty_cache()
            gc.collect()

    except Exception as e:
        print("Bỏ qua OCR do thiếu thư viện")
        
    return "\n".join(text)

def extract_text_docx(filepath):
    import docx
    import torch
    text = []
    try:
        doc = docx.Document(filepath)
        for para in doc.paragraphs:
            if para.text:
                text.append(para.text)
    except Exception as e:
        pass

    try:
        with zipfile.ZipFile(filepath, 'r') as archive:
            for item in archive.namelist():
                if item.startswith('word/media/'):
                    try:
                        image_data = archive.read(item)
                        if ocr_reader:
                            try:
                                ocr_result = ocr_reader.readtext(image_data, detail=0)
                                if ocr_result:
                                    text.append("\n".join(ocr_result))
                            except Exception as e:
                                print("Bỏ qua OCR do thiếu thư viện")
                        
                        del image_data
                        if use_gpu:
                            torch.cuda.empty_cache()
                        gc.collect()
                        
                    except Exception as e:
                        pass
    except Exception as e:
        pass
        
    return "\n".join(text)

def extract_text_pptx(filepath):
    from pptx import Presentation
    text = []
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
    except Exception as e:
        pass
    return "\n".join(text)

def extract_text_excel(filepath):
    import openpyxl
    text = []
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        text.append(str(cell))
    except Exception as e:
        pass
    return " ".join(text)

# ==========================================
# PHÂN LOẠI CHỦ ĐỀ AI (GPU)
# ==========================================
def classify_text(text):
    import torch
    if not classifier or not text.strip():
        return "Khác"
        
    try:
        # Cắt text thành các chunks nhỏ (khoảng 1000 ký tự một chunk)
        chunk_size = 1000
        text_chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
        
        # Chỉ lấy tối đa 8 chunks đầu tiên để tăng tốc độ phân loại
        text_chunks = text_chunks[:8]
        
        if not text_chunks:
            return "Khác"

        # Phân loại theo batch (tối ưu tốc độ) -> batch_size=8
        results = classifier(text_chunks, CANDIDATE_LABELS, batch_size=8)
        
        topic_scores = {label: 0 for label in CANDIDATE_LABELS}
        
        if isinstance(results, dict):
            results = [results]
            
        for res in results:
            for label, score in zip(res['labels'], res['scores']):
                topic_scores[label] += score
                
        best_topic = max(topic_scores, key=topic_scores.get)
        return best_topic
        
    except torch.cuda.OutOfMemoryError:
        print("Lỗi tràn VRAM GPU (OOM)! Đang tự động giải phóng bộ nhớ.")
        if use_gpu:
            torch.cuda.empty_cache()
        return "Khác"
    except Exception as e:
        print(f"Lỗi AI phân loại: {e}")
        return "Khác"

# ==========================================
# GIAO DIỆN CHÍNH
# ==========================================
class App:
    def __init__(self, root):
        self.root = root
        self.root.title("Desktop File Searcher (AI Indexing)")
        self.root.geometry("850x650")
        
        self.is_indexing = False
        init_db() # Khởi tạo DB nếu chưa có
        
        self.notebook = ttk.Notebook(root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.tab_index = ttk.Frame(self.notebook)
        self.notebook.add(self.tab_index, text="Quản lý Dữ liệu & Lập chỉ mục")
        self.setup_tab_index()
        
        self.tab_search = ttk.Frame(self.notebook)
        self.notebook.add(self.tab_search, text="Tìm Kiếm Nhanh")
        self.setup_tab_search()
        
        # System Tray setup
        self.root.protocol("WM_DELETE_WINDOW", self.hide_window)
        self.setup_tray()
        
        # Watchdog setup
        self.observer = None
        self.start_watchdog(self.entry_dir.get().strip())
        
    def setup_tab_index(self):
        frame_top = tk.Frame(self.tab_index)
        frame_top.pack(fill=tk.X, padx=10, pady=20)
        
        tk.Label(frame_top, text="Thư mục/Ổ đĩa:").grid(row=0, column=0, sticky=tk.W, pady=5)
        self.entry_dir = tk.Entry(frame_top, width=60)
        self.entry_dir.grid(row=0, column=1, padx=5, pady=5)
        self.entry_dir.insert(0, os.path.expanduser('~') if sys.platform == 'darwin' else "D:\\")
        
        btn_browse = tk.Button(frame_top, text="Chọn Thư Mục", command=self.browse_dir)
        btn_browse.grid(row=0, column=2, padx=5, pady=5)
        
        self.btn_index = tk.Button(frame_top, text="Bắt đầu Quét & Phân loại", command=self.start_indexing, bg="lightgreen")
        self.btn_index.grid(row=1, column=0, columnspan=3, pady=15)
        
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(frame_top, variable=self.progress_var, maximum=100)
        self.progress_bar.grid(row=2, column=0, columnspan=3, sticky=tk.EW, pady=10)
        
        self.lbl_index_status = tk.Label(frame_top, text="Trạng thái: Sẵn sàng.", fg="blue")
        self.lbl_index_status.grid(row=3, column=0, columnspan=3, sticky=tk.W, pady=5)
        
    def setup_tab_search(self):
        frame_top = tk.Frame(self.tab_search)
        frame_top.pack(fill=tk.X, padx=10, pady=10)
        
        tk.Label(frame_top, text="Từ khóa:").grid(row=0, column=0, sticky=tk.W, pady=5)
        self.entry_keyword = tk.Entry(frame_top, width=40)
        self.entry_keyword.grid(row=0, column=1, padx=5, pady=5)
        
        tk.Label(frame_top, text="Chủ đề:").grid(row=0, column=2, sticky=tk.W, pady=5)
        self.combo_topic = ttk.Combobox(frame_top, values=["Tất cả", "Toán học", "Lập trình", "Kiến thức Đại cương", "Khác"], state="readonly", width=20)
        self.combo_topic.current(0)
        self.combo_topic.grid(row=0, column=3, padx=5, pady=5)
        
        self.btn_search = tk.Button(frame_top, text="Tìm kiếm Siêu Tốc", command=self.do_search, bg="lightblue")
        self.btn_search.grid(row=0, column=4, padx=10, pady=5)
        
        self.lbl_search_status = tk.Label(frame_top, text="Nhập từ khóa và bấm Tìm kiếm.", fg="blue")
        self.lbl_search_status.grid(row=1, column=0, columnspan=5, sticky=tk.W, pady=5)
        
        self.paned_window = tk.PanedWindow(self.tab_search, orient=tk.VERTICAL)
        self.paned_window.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        frame_list = tk.Frame(self.paned_window)
        self.paned_window.add(frame_list, minsize=150)
        
        self.tree = ttk.Treeview(frame_list, columns=("Path", "Snippet", "Topic"), show="headings")
        self.tree.heading("Path", text="Đường dẫn File (Nhấp đúp để mở)")
        self.tree.heading("Snippet", text="Trích đoạn chứa từ khóa")
        self.tree.heading("Topic", text="Chủ đề")
        self.tree.column("Path", width=300)
        self.tree.column("Snippet", width=400)
        self.tree.column("Topic", width=100)
        self.tree.pack(fill=tk.BOTH, expand=True, side=tk.LEFT)
        
        scrollbar_list = ttk.Scrollbar(frame_list, orient=tk.VERTICAL, command=self.tree.yview)
        self.tree.configure(yscroll=scrollbar_list.set)
        scrollbar_list.pack(side=tk.RIGHT, fill=tk.Y)
        
        self.tree.bind("<Double-1>", self.open_file)
        self.tree.bind("<<TreeviewSelect>>", self.show_preview)
        
        frame_preview = tk.Frame(self.paned_window)
        self.paned_window.add(frame_preview, minsize=150)
        
        tk.Label(frame_preview, text="Xem trước nội dung (Từ khóa được bôi vàng):").pack(anchor=tk.W)
        self.text_preview = tk.Text(frame_preview, wrap=tk.WORD, state=tk.DISABLED)
        self.text_preview.pack(fill=tk.BOTH, expand=True, side=tk.LEFT)
        
        scrollbar_preview = ttk.Scrollbar(frame_preview, orient=tk.VERTICAL, command=self.text_preview.yview)
        self.text_preview.configure(yscroll=scrollbar_preview.set)
        scrollbar_preview.pack(side=tk.RIGHT, fill=tk.Y)
        
        self.text_preview.tag_config("highlight", background="yellow", foreground="black")
        
        self.current_search_results = {}
        self.current_search_keyword = ""

    def hide_window(self):
        self.root.withdraw()
        
    def show_window(self, icon=None, item=None):
        self.root.after(0, self.root.deiconify)
        
    def quit_app(self, icon, item):
        if getattr(self, 'observer', None):
            self.observer.stop()
            self.observer.join()
        icon.stop()
        self.root.after(0, self.root.destroy)

    def setup_tray(self):
        # Tạo icon ảnh đơn giản
        image = Image.new('RGB', (64, 64), color=(255, 0, 0))
        d = ImageDraw.Draw(image)
        d.rectangle((16, 16, 48, 48), fill=(0, 0, 255))
        
        menu = pystray.Menu(
            pystray.MenuItem("Mở giao diện", self.show_window, default=True),
            pystray.MenuItem("Thoát", self.quit_app)
        )
        self.tray_icon = pystray.Icon("DFS_AI", image, "Desktop Searcher AI", menu)
        threading.Thread(target=self.tray_icon.run, daemon=True).start()

    def start_watchdog(self, directory):
        if getattr(self, 'observer', None):
            self.observer.stop()
            self.observer.join()
            self.observer = None
            
        if directory and os.path.exists(directory):
            try:
                self.observer = Observer()
                handler = DocumentEventHandler(self)
                self.observer.schedule(handler, directory, recursive=True)
                self.observer.start()
            except Exception as e:
                print(f"Lỗi khởi tạo Watchdog: {e}")
        
    def browse_dir(self):
        d = filedialog.askdirectory()
        if d:
            self.entry_dir.delete(0, tk.END)
            self.entry_dir.insert(0, d)
            self.start_watchdog(d)
            
    def open_file(self, event):
        selected = self.tree.selection()
        if selected:
            filepath = self.tree.item(selected[0], "values")[0]
            try:
                if sys.platform == 'darwin':
                    subprocess.call(['open', filepath])
                else:
                    os.startfile(filepath)
            except Exception as e:
                messagebox.showerror("Lỗi", f"Không thể mở file: {e}")

    def show_preview(self, event):
        selected = self.tree.selection()
        if not selected:
            return
            
        filepath = self.tree.item(selected[0], "values")[0]
        content = self.current_search_results.get(filepath, "")
        keyword = self.current_search_keyword
        
        self.text_preview.config(state=tk.NORMAL)
        self.text_preview.delete(1.0, tk.END)
        self.text_preview.insert(tk.END, content)
        
        if keyword:
            idx = "1.0"
            first_match = None
            while True:
                idx = self.text_preview.search(keyword, idx, nocase=True, stopindex=tk.END)
                if not idx:
                    break
                if not first_match:
                    first_match = idx
                end_idx = f"{idx}+{len(keyword)}c"
                self.text_preview.tag_add("highlight", idx, end_idx)
                idx = end_idx
                
            if first_match:
                self.text_preview.see(first_match)
                
        self.text_preview.config(state=tk.DISABLED)

    def start_indexing(self):
        if self.is_indexing:
            return
            
        directory = self.entry_dir.get().strip()
        if not directory or not os.path.exists(directory):
            messagebox.showwarning("Cảnh báo", "Đường dẫn không hợp lệ!")
            return
            
        self.start_watchdog(directory)
            
        self.is_indexing = True
        self.btn_index.config(state=tk.DISABLED)
        self.lbl_index_status.config(text="Trạng thái: Đang quét siêu tốc metadata...")
        self.progress_var.set(0)
        
        threading.Thread(target=self.indexing_worker, args=(directory,), daemon=True).start()
        
    def indexing_worker(self, directory):
        allowed_extensions = {".pdf", ".docx", ".pptx", ".xlsx"}
        files_to_process = []
        
        for root_dir, _, files in os.walk(directory):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in allowed_extensions:
                    files_to_process.append(os.path.join(root_dir, file))
                    
        total_files = len(files_to_process)
        if total_files == 0:
            self.root.after(0, self.phase1_complete, 0, 0)
            return
            
        processed_count = 0
        skipped_count = 0
            
        for i, filepath in enumerate(files_to_process):
            try:
                os_mtime = os.path.getmtime(filepath)
                db_mtime = get_db_mtime(filepath)
                
                percent = ((i + 1) / total_files) * 100
                
                if os_mtime > db_mtime:
                    self.root.after(0, self.update_progress, percent, filepath, "Ghi nhận Metadata")
                    insert_or_update_document(filepath, "", "Đang chờ", os_mtime)
                    processed_count += 1
                else:
                    self.root.after(0, self.update_progress, percent, filepath, "Bỏ qua Metadata")
                    skipped_count += 1
                
            except Exception as e:
                print(f"Lỗi index metadata file {filepath}: {e}")
                
        self.root.after(0, self.phase1_complete, processed_count, skipped_count)
        
    def phase1_complete(self, processed_count, skipped_count):
        self.btn_index.config(state=tk.NORMAL)
        self.lbl_index_status.config(text=f"Hoàn tất quét siêu tốc! Đang khởi động đọc ngầm...")
        messagebox.showinfo("Thành công", f"Quá trình nạp Metadata hoàn tất.\n\n- File mới/cập nhật: {processed_count}\n- File đã bỏ qua: {skipped_count}\n\nĐã mở khóa tìm kiếm theo tên. Hệ thống đang tiến hành đọc nội dung ngầm.")
        threading.Thread(target=self.background_processing_worker, daemon=True).start()

    def background_processing_worker(self):
        try:
            conn = sqlite3.connect(DB_PATH, timeout=10)
            cursor = conn.cursor()
            cursor.execute("SELECT filepath FROM documents WHERE topic = 'Đang chờ'")
            pending_files = [row[0] for row in cursor.fetchall()]
            conn.close()
        except Exception as e:
            print(f"Lỗi lấy danh sách phase 2: {e}")
            self.is_indexing = False
            return
            
        if not pending_files:
            self.root.after(0, lambda: self.lbl_index_status.config(text="Trạng thái: Hoàn tất 100% (Không có file cần đọc nội dung)."))
            self.is_indexing = False
            return
            
        def status_callback(msg):
            self.root.after(0, lambda: self.lbl_index_status.config(text=msg))
        init_ai_models(status_callback)
        
        total = len(pending_files)
        for i, filepath in enumerate(pending_files):
            percent = ((i + 1) / total) * 100
            self.root.after(0, self.update_progress, percent, filepath, "Phân tích ngầm nội dung")
            
            ext = os.path.splitext(filepath)[1].lower()
            try:
                os_mtime = os.path.getmtime(filepath)
                content = ""
                if ext == ".pdf":
                    content = extract_text_pdf(filepath)
                elif ext == ".docx":
                    content = extract_text_docx(filepath)
                elif ext == ".pptx":
                    content = extract_text_pptx(filepath)
                elif ext == ".xlsx":
                    content = extract_text_excel(filepath)
                    
                topic = classify_text(content)
                insert_or_update_document(filepath, content, topic, os_mtime)
                
            except Exception as e:
                print(f"Lỗi đọc nội dung file {filepath}: {e}")
                insert_or_update_document(filepath, "", "Lỗi đọc", get_db_mtime(filepath))
                
            time.sleep(0.1) # Throttling nhường CPU/GPU
            
        self.root.after(0, lambda: self.lbl_index_status.config(text="Trạng thái: Đã hoàn tất 100% lập chỉ mục toàn bộ nội dung."))
        self.is_indexing = False

    def update_progress(self, percent, current_file, action_text):
        self.progress_var.set(percent)
        self.lbl_index_status.config(text=f"Trạng thái: {action_text} ({percent:.1f}%) - {os.path.basename(current_file)}")

    def do_search(self):
        keyword = self.entry_keyword.get().strip().lower()
        topic = self.combo_topic.get()
        
        if not keyword:
            messagebox.showwarning("Cảnh báo", "Vui lòng nhập từ khóa tìm kiếm!")
            return
            
        # Xóa kết quả cũ
        for item in self.tree.get_children():
            self.tree.delete(item)
            
        self.text_preview.config(state=tk.NORMAL)
        self.text_preview.delete(1.0, tk.END)
        self.text_preview.config(state=tk.DISABLED)
        
        self.current_search_results.clear()
        self.current_search_keyword = keyword
            
        self.lbl_search_status.config(text="Đang tìm kiếm trong Database...")
        self.root.update()
        
        results = search_db(keyword, topic)
        
        for filepath, content, item_topic in results:
            self.current_search_results[filepath] = content
            
            content_lower = content.lower()
            idx = content_lower.find(keyword)
            if idx != -1:
                start_idx = max(0, idx - 40)
                end_idx = min(len(content), idx + len(keyword) + 100)
                snippet = content[start_idx:end_idx].replace("\n", " ")
                if start_idx > 0:
                    snippet = "..." + snippet
                if end_idx < len(content):
                    snippet = snippet + "..."
            else:
                snippet = "Không tìm thấy đoạn chứa từ khóa trực tiếp."
                
            self.tree.insert("", tk.END, values=(filepath, snippet, item_topic))
            
        self.lbl_search_status.config(text=f"Tìm thấy {len(results)} kết quả.")

if __name__ == "__main__":
    root = tk.Tk()
    app = App(root)
    root.mainloop()
